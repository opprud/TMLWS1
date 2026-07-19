#!/usr/bin/env python3
"""Inject Marp speaker notes (HTML comments) into an editable PPTX.

marp-cli's experimental --pptx-editable drops speaker notes; this puts them
back. The first plain HTML comment of each slide's markdown (ignoring Marp
directives like `_class:`) becomes the slide's presenter note.

Usage: python3 inject_notes.py <deck.md> <deck.pptx>
"""
import re
import sys

from pptx import Presentation


def md_notes(md_path):
    """Return the speaker note (or '') for each slide in the deck."""
    src = open(md_path, encoding="utf8").read()
    # strip front-matter, then split on slide separators
    body = re.sub(r"\A---\n.*?\n---\n", "", src, count=1, flags=re.S)
    slides = re.split(r"\n---\s*\n", body)
    notes = []
    for slide in slides:
        note_parts = []
        for m in re.finditer(r"<!--(.*?)-->", slide, flags=re.S):
            text = m.group(1).strip()
            if re.match(r"_?[a-zA-Z-]+\s*:", text) and "\n" not in text:
                continue  # Marp directive (e.g. _class: title), not a note
            if text.upper().startswith("TODO"):
                continue
            note_parts.append(text)
        notes.append("\n\n".join(note_parts))
    return notes


def main():
    md_path, pptx_path = sys.argv[1], sys.argv[2]
    notes = md_notes(md_path)
    prs = Presentation(pptx_path)
    if len(prs.slides) != len(notes):
        print(f"WARNING: {len(prs.slides)} pptx slides vs {len(notes)} md slides — matching by index")
    injected = 0
    for slide, note in zip(prs.slides, notes):
        if note:
            slide.notes_slide.notes_text_frame.text = note
            injected += 1
    prs.save(pptx_path)
    print(f"{pptx_path}: injected notes into {injected}/{len(prs.slides)} slides")


if __name__ == "__main__":
    main()
